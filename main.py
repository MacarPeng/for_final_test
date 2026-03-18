import os
import sys

# --- 针对 D 盘路径的 GPU 补丁自动修复 ---
if os.name == 'nt':
    # 获取当前 venv 的 site-packages 路径
    venv_path = os.path.join(os.getcwd(), 'venv', 'Lib', 'site-packages')
    nvidia_dir = os.path.join(venv_path, 'nvidia')
    
    if os.path.exists(nvidia_dir):
        # 遍历所有可能的 dll 目录
        for root, dirs, files in os.walk(nvidia_dir):
            if 'bin' in dirs:
                bin_path = os.path.join(root, 'bin')
                os.add_dll_directory(bin_path)
                print(f"✅ 已手动加载 GPU 依赖库: {bin_path}")
# ---------------------------------------

from moviepy import VideoFileClip
from pptx import Presentation
from PIL import Image
import io
from faster_whisper import WhisperModel
import jieba
from collections import Counter
import cv2
from fpdf import FPDF
from datetime import datetime
import jieba.analyse 
from tqdm import tqdm
from openai import OpenAI
import json
from prompt_config import SKILLS, PPT_SELECTION_PROMPT, GLOBAL_SUMMARY_PROMPT
from keywords_config import DEFAULT_KEYWORDS, COURSE_KEYWORDS
from stop_words import STOP_WORDS

# 语音识别模型配置
# 可选: "small", "medium", "large-v3" 
MODEL_SIZE = "small"
USE_CPU = False  # True 强制使用CPU，False 使用GPU

# 目录配置
VIDEO_DIR = "videos"
AUDIO_DIR = "audios"
SCREENSHOT_DIR = "screenshots"
RESULT_DIR = "results"
PPT_DIR = "pptx"

FONT_PATH = r"C:\Windows\Fonts\simhei.ttf"

def clean_text_for_pdf(text):
    """清理文本中可能导致PDF错误的字符"""
    if not text:
        return ""
    # 移除或替换可能导致问题的字符
    import re
    # 替换一些特殊字符
    text = text.replace('\x00', '')
    text = text.replace('\ufeff', '')
    # 移除emoji等特殊符号（可选）
    # text = re.sub(r'[\U00010000-\U0010ffff]', '', text)
    return text

def clean_cache():
    """清理可能的损坏缓存"""
    import shutil
    cache_dirs = [
        os.path.expanduser("~/.cache/huggingface"),
        os.path.join(os.getcwd(), "cache"),
    ]
    for cache_dir in cache_dirs:
        if os.path.exists(cache_dir):
            try:
                shutil.rmtree(cache_dir)
                print(f"✅ 已清理缓存: {cache_dir}")
            except:
                pass

def extract_audio(video_path,output_audio):
    video=VideoFileClip(video_path)
    audio=video.audio
    if audio is None:
        print(f"❌ 视频没有音频轨道: {video_path}")
        video.close()
        return
    video.audio.write_audiofile(output_audio,fps=16000,nbytes=2,codec='pcm_s16le')
    video.close()

def save_ppt_screenshot(v_path, time_s, out_name):
    """
    使用 MoviePy 在指定时间点截取一帧
    """
    try:
        # 确保保存图片的文件夹存在
        os.makedirs(os.path.dirname(out_name), exist_ok=True)
        
        # 加载视频并截图
        with VideoFileClip(v_path) as clip:
            # 这里的 time_s 是秒，clip.save_frame 会自动定位
            # 稍微往后偏 0.5 秒，避开可能存在的转场黑屏
            target_t = min(time_s + 0.5, clip.duration - 0.1)
            clip.save_frame(out_name, t=target_t)
            
        print(f"📸 截图成功: {out_name}")
        return True
    except Exception as e:
        print(f"❌ 截图失败: {e}")
        return False

def transcribe_audio(model_obj, audio_path):
    segments_generator, info = model_obj.transcribe(audio_path, beam_size=5, language="zh")
    segments_list = []
    print(f"--- 音频总时长: {info.duration:.1f} 秒 ---")
    with tqdm(total=round(info.duration, 1), unit="s", desc="语音识别") as pbar:
        last_pos = 0
        for s in segments_generator:
            segments_list.append({
                'start': s.start,
                'end': s.end,
                'text': s.text
            })
            tqdm.write(f"[{s.start:.1f}s -> {s.end:.1f}s]: {s.text}")
            increment = s.end - last_pos
            if increment > 0:
                pbar.update(increment)
                last_pos = s.end
                
    return segments_list

def find_keyword_hits(segments, keywords):
    hit_times = []
    last_hit_time = -999
    for s in segments:
        text = s['text']
        start_time = s['start']
        found = any(kw in text for kw in keywords)
        if found:
            if start_time - last_hit_time > 120:  
                hit_times.append(start_time)
                last_hit_time = start_time
                print(f"🚩 命中关键词，时间点：{start_time:.1f}s，内容：{text}")
    return hit_times

def analyze_and_update_stop_words(all_segments):
    """利用LLM分析课程词频，找出高频无意义词汇"""
    print("\n--- 正在分析课程词频，生成个性化stop_words ---")
    
    # 1. 先用jieba分词并统计词频
    all_text = "".join([s['text'] for s in all_segments])
    words = jieba.lcut(all_text)
    
    # 过滤单字和数字
    words = [w for w in words if len(w) > 1 and not w.isdigit()]
    word_counts = Counter(words)
    
    # 获取高频词（前100个）
    top_words = word_counts.most_common(100)
    top_words_str = "\n".join([f"{w}: {c}次" for w, c in top_words[:50]])
    
    # 2. 让LLM从中筛选无意义词汇
    prompt = f"""
以下是一段课堂录音中按出现频率排序的词语（词:次数）：

{top_words_str}

请找出其中出现频率很高但对学习无意义的词汇（如口语词、虚词、过渡词、重复词等）。
这些词应该被加入stop_words列表中过滤掉。

注意：专业术语、概念名词即使频率高也要保留。

请返回JSON格式：
{{
    "meaningless_words": ["词1", "词2", ...],
    "reason": "简要说明每个词为什么无意义"
}}
"""
    try:
        response = client.chat.completions.create(
            model="deepseek-chat",
            messages=[
                {"role": "system", "content": "你是一个文本分析专家，擅长识别无意义的词汇。"},
                {"role": "user", "content": prompt},
            ],
            response_format={'type': 'json_object'},
            stream=False
        )
        content_str = response.choices[0].message.content
        res_dict = json.loads(content_str)
        
        new_stop_words = res_dict.get("meaningless_words", [])
        reason = res_dict.get("reason", "")
        
        if new_stop_words:
            print(f"发现需要过滤的高频词: {new_stop_words}")
            print(f"原因: {reason}")
            
            # 读取现有stop_words
            with open("stop_words.py", "r", encoding="utf-8") as f:
                content = f.read()
            
            # 提取现有词汇列表
            import re
            match = re.search(r'STOP_WORDS\s*=\s*\[(.*?)\]', content, re.DOTALL)
            if match:
                existing_words = re.findall(r'"([^"]+)"', match.group(1))
                # 添加新词汇（去重）
                updated_words = list(set(existing_words + new_stop_words))
                # 按字母排序
                updated_words.sort()
                
                # 重新生成文件内容
                new_content = '# Stop Words - 常见无意义词汇（口语、虚词等）\n'
                new_content += '# 这些词会在关键词提取时被过滤掉\n\n'
                new_content += 'STOP_WORDS = [\n'
                for w in updated_words:
                    new_content += f'    "{w}",\n'
                new_content += ']\n'
                
                with open("stop_words.py", "w", encoding="utf-8") as f:
                    f.write(new_content)
                
                print(f"✅ 已更新 stop_words.py，当前共 {len(updated_words)} 个词汇")
                return new_stop_words
            else:
                print("❌ 无法解析 stop_words.py 格式")
                return []
        else:
            print("未发现需要过滤的高频词")
            return []
            
    except Exception as e:
        print(f"❌ 分析失败: {e}")
        return []

def generate_course_keywords(all_segments, course_name=""):
    """利用LLM根据课程名和内容自动生成重点关键词"""
    print("\n--- 正在生成课程专属关键词 ---")
    
    # 取前3000字符作为样本
    sample_text = "".join([s['text'] for s in all_segments[:50]])[:3000]
    
    prompt = f"""
根据以下课程信息，自动推断并生成适合这门课的"考点关键词"列表。

课程名称: {course_name or "未知（请根据内容推断）"}
课堂前部分内容:
{sample_text}

要求：
1. 先推断这是什么课程（根据名称和内容）
2. 生成15-20个这门课很可能考到的关键词
3. 关键词应该能帮助学生在课堂录音中找到考点

请返回JSON格式：
{{
    "inferred_course": "推断的课程名",
    "course_keywords": ["关键词1", "关键词2", ...],
    "reason": "为什么这些关键词可能是考点"
}}
"""
    try:
        response = client.chat.completions.create(
            model="deepseek-chat",
            messages=[
                {"role": "system", "content": "你是一个资深大学教授，擅长预测课程考点。"},
                {"role": "user", "content": prompt},
            ],
            response_format={'type': 'json_object'},
            stream=False
        )
        content_str = response.choices[0].message.content
        res_dict = json.loads(content_str)
        
        inferred_course = res_dict.get("inferred_course", "")
        course_keywords = res_dict.get("course_keywords", [])
        reason = res_dict.get("reason", "")
        
        if course_keywords:
            print(f"📚 推断课程: {inferred_course}")
            print(f"📝 生成关键词: {course_keywords}")
            print(f"💡 原因: {reason}")
            return course_keywords
        else:
            print("未能生成课程关键词")
            return []
            
    except Exception as e:
        print(f"❌ 生成失败: {e}")
        return []

def analyze_hits_and_extract_ppt(hit_times, all_segments, save_dir, ppt_path=None, slides_text=None):
    """
    分析命中点并从PDF PPT中提取相关内容
    不再从视频截图，改为从PPT PDF中选择页面
    """
    results = []
    
    # 如果有PPT，预加载文字内容
    if ppt_path and slides_text:
        print(f"📚 已加载PPT，共 {len(slides_text)} 页")
    
    for i, t in enumerate(hit_times):
        start_range, end_range = max(0, t - 60), t + 60
        
        # 获取这个窗口内的所有 segment
        window_segments = [s for s in all_segments if start_range <= s['start'] <= end_range]
        window_text = "".join([s['text'] for s in window_segments])

        # 请求 DeepSeek 分析
        print(f"🧠 正在分析第 {i+1}/{len(hit_times)} 个重点片段...")
        ai_data = get_llm_advanced_analysis(window_segments, t)
        
        # 检查内容是否相关，不相关则跳过
        is_relevant = ai_data.get('is_relevant', True)
        if is_relevant is False or is_relevant == "false":
            print(f"   ⏭️ 内容不相关，跳过")
            continue
        
        best_t = ai_data.get('best_time', t)
        
        # 从PPT中选择相关页面
        related_ppt_pages = []
        if ppt_path and slides_text:
            # 构造这个片段的摘要给LLM
            segment_info = {
                'time': best_t,
                'summary': ai_data.get('summary', ''),
                'keywords': ai_data.get('keywords', []),
                'exam_points': ai_data.get('exam_points', [])
            }
            related_pages = find_related_ppt_pages(slides_text, segment_info, window_text)
            related_ppt_pages = related_pages
        
        results.append({
            'hit_time': best_t,
            'top_words': ai_data.get('keywords', []),
            'summary': ai_data.get('summary', ''),
            'exam_points': ai_data.get('exam_points', []),
            'review_tips': ai_data.get('review_tips', ''),
            'structured_points': ai_data.get('structured_points', []),  # 结构化重点
            'related_ppt_pages': related_ppt_pages  # 关联的PPT页面
        })
    
    return results

def find_related_ppt_pages(slides_text, segment_info, window_text):
    """根据片段信息从PPT中找到相关页面"""
    prompt = f"""
你是一个学术助手。现在需要从PPT中找出与课堂重点片段相关的页面。

【课堂重点片段信息】
- 时间点: {segment_info['time']:.1f}秒
- 关键词: {", ".join(segment_info['keywords'])}
- 考点: {", ".join(segment_info['exam_points'])}
- 摘要: {segment_info['summary']}

【课堂原文片段】（用于理解上下文）:
{window_text[:1000]}

【PPT各页内容】
{chr(10).join([f"第{s['page']}页: {s['text'][:300]}..." for s in slides_text[:30]])}

请找出与这个课堂片段最相关的PPT页面（可能1-2页）。

要求：
1. 只选择真正讲解了关键词/考点的页面
2. 不要选择无关页面（如目录页、封面页等，除非确实相关）
3. 返回页码列表

JSON格式：
{{"related_pages": [页码1], "reason": "理由"}}
"""
    
    try:
        response = client.chat.completions.create(
            model="deepseek-chat",
            messages=[
                {"role": "system", "content": "你是一个严谨的学术助手，擅长匹配课堂重点和PPT内容。"},
                {"role": "user", "content": prompt},
            ],
            response_format={'type': 'json_object'},
            stream=False
        )
        content_str = response.choices[0].message.content
        res_dict = json.loads(content_str)
        pages = res_dict.get("related_pages", [])
        reason = res_dict.get("reason", "")
        
        if pages:
            print(f"   📄 关联PPT页面: {pages} - {reason}")
        return pages
    except Exception as e:
        print(f"   ⚠️ PPT关联失败: {e}")
        return []

def extract_ppt_text(ppt_path):
    """提取PPT所有页面的文字（支持pptx和pdf）"""
    if ppt_path.endswith('.pdf'):
        return extract_pdf_text(ppt_path)
    else:
        prs = Presentation(ppt_path)
        slides_text = []
        for i, slide in enumerate(prs.slides):
            text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + " "
            slides_text.append({
                "page": i + 1,
                "text": text.strip()
            })
        return slides_text

def extract_pdf_text(pdf_path):
    """提取PDF所有页面的文字"""
    from pypdf import PdfReader
    reader = PdfReader(pdf_path)
    slides_text = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        slides_text.append({
            "page": i + 1,
            "text": text.strip()
        })
    return slides_text

def select_pdf_pages(pdf_path, selected_pages, output_path):
    """从PDF中提取指定页面生成新PDF"""
    from pypdf import PdfReader, PdfWriter
    reader = PdfReader(pdf_path)
    writer = PdfWriter()
    
    for page_num in selected_pages:
        if 0 <= page_num < len(reader.pages):
            writer.add_page(reader.pages[page_num])
    
    with open(output_path, 'wb') as f:
        writer.write(f)
    print(f"📄 已提取PDF页面: {selected_pages}")

def select_relevant_slides(slides_text, keywords, hit_times_text):
    """使用LLM根据关键词和课堂重点筛选相关PPT页面"""
    slides_summary = "\n".join([
        f"第{s['page']}页: {s['text'][:200]}..."
        for s in slides_text[:30]  # 限制数量避免超过token限制
    ])
    
    prompt = f"""
你是一个学术助手。我会提供：
1. 课堂中老师讲到的重点内容（可能包含多个片段）：
{hit_times_text}
2. PPT所有页面的文字内容：
{slides_summary}

请根据课堂重点内容，从PPT中选出最相关的页面（可能1-3页）。

请严格按以下JSON格式回复：
{{
    "selected_pages": [页码1, 页码2],
    "reason": "简要说明为什么选这些页面"
}}
"""
    
    try:
        response = client.chat.completions.create(
            model="deepseek-chat",
            messages=[
                {"role": "system", "content": "你是一个严谨的学术助手，擅长匹配课堂重点和PPT内容。"},
                {"role": "user", "content": prompt},
            ],
            response_format={'type': 'json_object'},
            stream=False
        )
        content_str = response.choices[0].message.content
        res_dict = json.loads(content_str)
        return res_dict.get("selected_pages", []), res_dict.get("reason", "")
    except Exception as e:
        print(f"❌ LLM筛选PPT失败: {e}")
        return [], ""

def convert_ppt_to_images(ppt_path, selected_pages, output_dir):
    """将PPT指定页面转换为图片"""
    prs = Presentation(ppt_path)
    saved_images = []
    
    for page_num in selected_pages:
        if 0 < page_num <= len(prs.slides):
            slide = prs.slides[page_num - 1]
            
            # 创建对应尺寸的图片
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            # 使用PIL创建空白图片
            from PIL import ImageDraw, ImageFont
            dpi = 150
            width_px = int(slide_width * dpi / 914400)
            height_px = int(slide_height * dpi / 914400)
            
            img = Image.new('RGB', (width_px, height_px), (255, 255, 255))
            draw = ImageDraw.Draw(img)
            
            # 提取文字绘制到图片上（简化版：纯文字PPT）
            y_offset = 20
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    try:
                        draw.text((20, y_offset), shape.text[:100], fill=(0, 0, 0))
                        y_offset += 30
                    except:
                        pass
            
            img_path = os.path.join(output_dir, f"ppt_page_{page_num}.jpg")
            img.save(img_path, "JPEG", quality=95)
            saved_images.append({
                "page": page_num,
                "img_path": img_path
            })
            print(f"📄 已转换PPT第{page_num}页为图片")
    
    return saved_images

# ================= 彻底修复版 PDF 模块 =================

class MyCoursePDF(FPDF):
    def footer(self):
        """底栏页码"""
        self.set_y(-15)
        self.set_font("Hans", size=10)
        self.set_text_color(150, 150, 150)
        # 使用 fpdf2 推荐的 new_x 模式
        self.cell(0, 10, text=f"第 {self.page_no()} 页", align="C")

def create_final_pdf(analysis_results, course_name, filename):
    """创建课堂笔记PDF - 终极修复版"""
    print(f"📑 正在排版生成：{filename}")
    pdf = MyCoursePDF()
    
    # 注册字体
    try:
        pdf.add_font("Hans", "", FONT_PATH)
        pdf.add_font("Hans", "B", FONT_PATH)
    except:
        print("❌ 警告：字体加载失败")

    # --- 1. 封面 ---
    pdf.add_page()
    pdf.set_font("Hans", size=30)
    pdf.ln(60)
    pdf.cell(0, 20, text=course_name, align='C', new_x="LMARGIN", new_y="NEXT")
    pdf.set_font("Hans", size=18)
    pdf.cell(0, 20, text="AI 智能辅助复习讲义", align='C', new_x="LMARGIN", new_y="NEXT")
    pdf.ln(40)
    pdf.set_font("Hans", size=12)
    pdf.cell(0, 10, text=f"生成日期：{datetime.now().strftime('%Y-%m-%d')}", align='C', new_x="LMARGIN", new_y="NEXT")

    # --- 2. 内容页 ---
    for item in analysis_results:
        pdf.add_page()
        
        # A. 时间标题
        t = item.get('hit_time', 0)
        pdf.set_fill_color(240, 240, 240)
        pdf.set_font("Hans", size=15)
        time_text = f"重点片段：{int(t // 60)}分{int(t % 60)}秒"
        pdf.cell(180, 12, text=time_text, fill=True, new_x="LMARGIN", new_y="NEXT")
        pdf.ln(5)

        # B. 核心词汇 - 强制重置 X 坐标
        pdf.set_x(10)
        pdf.set_font("Hans", size=12)
        pdf.set_text_color(31, 73, 125)
        keywords = item.get('top_words', [])
        # 核心修复：将 0 宽度改为 180 固定宽度
        pdf.multi_cell(180, 10, text="核心词汇：" + " | ".join(keywords))
        pdf.ln(2)
        
        # C. 结构化重点 (Bullet Points) - 修复崩溃点
        if item.get('structured_points'):
            pdf.set_text_color(80, 80, 80)
            pdf.set_font("Hans", size=11, style="B")
            pdf.set_x(10)
            pdf.cell(180, 8, text="📋 重点详解：", new_x="LMARGIN", new_y="NEXT")
            
            for point in item['structured_points']:
                pdf.set_font("Hans", size=10, style="")
                pdf.set_x(15) # 设置缩进
                # 核心修复：使用固定宽度 175，避免 0 宽度计算错误
                pdf.multi_cell(175, 7, text=f"• {point}")
                pdf.ln(1)
            pdf.ln(3)

        # D. AI 摘要 - 修复颜色和背景冲突
        pdf.set_fill_color(252, 243, 207)
        pdf.set_text_color(0, 0, 0)
        pdf.set_font("Hans", size=12, style="B")
        pdf.set_x(10)
        pdf.cell(180, 8, text="【AI 考点总结】：", new_x="LMARGIN", new_y="NEXT")
        
        pdf.set_font("Hans", size=11, style="")
        pdf.set_x(10)
        summary_text = clean_text_for_pdf(item.get('summary', ''))
        # 核心修复：固定宽度 180
        pdf.multi_cell(180, 8, text=summary_text, fill=True)
        
        # E. 复习建议
        if item.get('review_tips'):
            pdf.ln(5)
            pdf.set_fill_color(230, 245, 230)
            pdf.set_text_color(0, 100, 0)
            pdf.set_x(10)
            review_text = clean_text_for_pdf(item['review_tips'])
            pdf.multi_cell(180, 7, text="📖 复习建议：" + review_text, fill=True)

    # --- 3. 保存 ---
    try:
        pdf.output(filename)
        print(f"✨ 最终讲义已成功保存至: {filename}")
    except Exception as e:
        print(f"❌ PDF 导出失败: {e}")

# 2. 如果你需要带 PPT 截图的 PDF，也请使用下面这个修复版
def create_pdf_with_ppt(analysis_results, ppt_images, course_name, filename):
    """创建包含 PPT 页面和 AI 课堂分析的综合 PDF - 深度排版加固版"""
    print(f"📑 正在排版生成综合讲义：{filename}")
    pdf = MyCoursePDF()
    
    # 1. 注册字体 (确保使用本地存在的路径)
    try:
        pdf.add_font("Hans", "", FONT_PATH)
        pdf.add_font("Hans", "B", FONT_PATH)
    except:
        print("❌ 警告：字体加载失败")

    # --- 第一部分：封面 ---
    pdf.add_page()
    pdf.set_font("Hans", size=30)
    pdf.ln(60)
    pdf.cell(0, 20, text=course_name, align='C', new_x="LMARGIN", new_y="NEXT")
    pdf.set_font("Hans", size=18)
    pdf.cell(0, 20, text="AI 智能辅助复习讲义 (含 PPT 重点页)", align='C', new_x="LMARGIN", new_y="NEXT")
    pdf.ln(40)
    pdf.set_font("Hans", size=12)
    pdf.cell(0, 10, text=f"生成日期：{datetime.now().strftime('%Y-%m-%d %H:%M')}", align='C', new_x="LMARGIN", new_y="NEXT")

    # --- 第二部分：PPT 重点页面汇总 (如果有) ---
    if ppt_images:
        pdf.add_page()
        pdf.set_font("Hans", size=22, style="B")
        pdf.set_text_color(44, 62, 80)
        pdf.cell(0, 20, text="📚 关联 PPT 重点页面", align='C', new_x="LMARGIN", new_y="NEXT")
        pdf.ln(10)
        
        for item in ppt_images:
            img_path = item['img_path']
            if os.path.exists(img_path):
                # 检查页面剩余空间，如果不够放图则换页
                if pdf.get_y() > 180:
                    pdf.add_page()
                
                pdf.image(img_path, x=15, w=180)
                pdf.set_y(pdf.get_y() + 5) # 稍微往下挪一点
                pdf.set_font("Hans", size=10)
                pdf.set_text_color(100, 100, 100)
                pdf.cell(0, 10, text=f"PPT 第 {item['page']} 页", align='C', new_x="LMARGIN", new_y="NEXT")
                pdf.ln(10)
        pdf.set_text_color(0, 0, 0) # 颜色恢复黑色

    # --- 第三部分：课堂重点深度分析 ---
    for item in analysis_results:
        pdf.add_page()
        
        # A. 标题栏
        t = item.get('hit_time', 0)
        pdf.set_fill_color(240, 240, 240)
        pdf.set_font("Hans", size=15, style="B")
        time_text = f"重点片段：{int(t // 60)}分{int(t % 60)}秒"
        pdf.cell(180, 12, text=time_text, fill=True, new_x="LMARGIN", new_y="NEXT")
        pdf.ln(5)

        # B. 核心词汇
        pdf.set_x(10) # 强制重置 X 坐标
        pdf.set_font("Hans", size=12)
        pdf.set_text_color(31, 73, 125)
        keywords = item.get('top_words', [])
        pdf.multi_cell(180, 10, text="核心词汇：" + " | ".join(keywords))
        pdf.ln(2)
        
        # C. 结构化重点 (详解列表)
        if item.get('structured_points'):
            pdf.set_text_color(80, 80, 80)
            pdf.set_font("Hans", size=11, style="B")
            pdf.set_x(10)
            pdf.cell(180, 8, text="📋 重点详解：", new_x="LMARGIN", new_y="NEXT")
            
            for point in item['structured_points']:
                pdf.set_font("Hans", size=10, style="")
                pdf.set_x(15) # 设置缩进
                # 使用固定宽度 175，防止 horizontal space 报错
                pdf.multi_cell(175, 7, text=f"• {point}")
            pdf.ln(3)

        # D. AI 考点总结 (黄色背景区域)
        pdf.set_fill_color(252, 243, 207)
        pdf.set_text_color(0, 0, 0)
        pdf.set_font("Hans", size=12, style="B")
        pdf.set_x(10)
        pdf.cell(180, 8, text="【AI 考点总结】：", new_x="LMARGIN", new_y="NEXT")
        
        pdf.set_font("Hans", size=11, style="")
        pdf.set_x(10)
        summary_text = clean_text_for_pdf(item.get('summary', ''))
        # 强制 180 宽度
        pdf.multi_cell(180, 8, text=summary_text, fill=True)
        
        # E. 复习建议 (绿色背景区域)
        if item.get('review_tips'):
            pdf.ln(5)
            pdf.set_fill_color(230, 245, 230)
            pdf.set_text_color(0, 100, 0)
            pdf.set_x(10)
            pdf.set_font("Hans", size=11, style="B")
            pdf.cell(180, 8, text="📖 复习建议：", new_x="LMARGIN", new_y="NEXT")
            pdf.set_font("Hans", size=10, style="")
            pdf.set_x(10)
            review_text = clean_text_for_pdf(item['review_tips'])
            pdf.multi_cell(180, 7, text=review_text, fill=True)

    # --- 最后一步：保存 PDF ---
    try:
        pdf.output(filename)
        print(f"✨ 包含 PPT 的综合讲义已成功保存至: {filename}")
    except Exception as e:
        print(f"❌ PDF 导出失败: {e}")

# 1. 初始化客户端 (确保 API Key 已填)
client = OpenAI(
    api_key="YOUR_API_KEY_HERE", 
    base_url="https://api.deepseek.com",
    timeout=60.0 # 长视频分析较慢，设置60秒超时
)

CURRENT_SKILL = "default"  # 可选: "default", "detailed", "simple"

def get_llm_advanced_analysis(window_segments, default_time):
    """
    使用配置好的skill来分析课堂内容
    """
    skill = SKILLS.get(CURRENT_SKILL, SKILLS["default"])
    
    context_with_time = ""
    for s in window_segments:
        context_with_time += f"[{s['start']:.1f}s] {s['text']}\n"

    prompt = skill["analysis_prompt"].format(context=context_with_time)
    system_prompt = skill["system_prompt"]

    try:
        response = client.chat.completions.create(
            model="deepseek-chat",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": prompt},
            ],
            response_format={'type': 'json_object'},
            stream=False
        )

        content_str = response.choices[0].message.content
        res_dict = json.loads(content_str)
        
        # 检查内容是否相关
        is_relevant = res_dict.get("is_relevant", True)
        
        # 如果不相关，返回特殊标记
        if is_relevant is False or is_relevant == "false":
            return {
                "is_relevant": False,
                "keywords": [],
                "summary": "",
                "exam_points": [],
                "review_tips": "",
                "structured_points": [],
                "best_time": default_time
            }
        
        keywords = res_dict.get("keywords", ["重点内容"])
        summary = res_dict.get("summary", "未能生成摘要")
        exam_points = res_dict.get("exam_points", [])
        review_tips = res_dict.get("review_tips", "")
        structured_points = res_dict.get("structured_points", [])  # 结构化重点
        
        ai_time = res_dict.get("best_time", default_time)
        
        if window_segments:
            if not (window_segments[0]['start'] - 5 <= ai_time <= window_segments[-1]['end'] + 5):
                print(f"⚠️ AI 建议的时间 {ai_time}s 超出窗口，使用默认值 {default_time}s")
                ai_time = default_time
        else:
            ai_time = default_time
            
        return {
            "is_relevant": True,
            "keywords": keywords,
            "summary": summary,
            "exam_points": exam_points,
            "review_tips": review_tips,
            "structured_points": structured_points,
            "best_time": float(ai_time)
        }

    except Exception as e:
        print(f"❌ DeepSeek 分析失败: {e}")
        return {
            "keywords": ["识别失败"],
            "summary": "由于网络或 API 原因，未能生成 AI 摘要。",
            "exam_points": [],
            "review_tips": "",
            "best_time": default_time
        }
# ================= 主函数 =================
if __name__ == "__main__":
    # 1. 自动创建基础文件夹
    for d in [VIDEO_DIR, AUDIO_DIR, SCREENSHOT_DIR, RESULT_DIR, PPT_DIR]:
        os.makedirs(d, exist_ok=True)
    video_files = [f for f in os.listdir(VIDEO_DIR) if f.endswith(('.mp4', '.mkv', '.avi','.ts'))]
    if not video_files:
        print(f"❌ '{VIDEO_DIR}' 文件夹里没有视频文件！")
        import sys
        sys.exit()
    print("--- 发现以下视频 ---")
    for i, f in enumerate(video_files):
        print(f"{i+1}. {f}")
    choice = int(input("请输入要处理的视频编号: ")) - 1
    target_video_name = video_files[choice]
    base_name = os.path.splitext(target_video_name)[0]
    VIDEO_PATH = os.path.join(VIDEO_DIR, target_video_name)
    AUDIO_PATH = os.path.join(AUDIO_DIR, f"{base_name}.wav")
    current_screenshot_dir = os.path.join(SCREENSHOT_DIR, base_name)
    os.makedirs(current_screenshot_dir, exist_ok=True)
    PDF_PATH = os.path.join(RESULT_DIR, f"{base_name}_复习讲义.pdf")
    
    # 选择分析模式
    print("\n--- 选择分析模式 ---")
    print("1. default - 标准模式（关键词+摘要+考点预测）")
    print("2. detailed - 详细模式（更全面的分析）")
    print("3. simple - 简洁模式（快速提取）")
    skill_choice = input("请输入模式编号（默认1）: ").strip() or "1"
    skill_map = {"1": "default", "2": "detailed", "3": "simple"}
    CURRENT_SKILL = skill_map.get(skill_choice, "default")
    print(f"已选择模式: {SKILLS[CURRENT_SKILL]['name']}")
    
    # PPT上传处理
    ppt_path = None
    ppt_images = []
    print("\n--- PPT 处理（可选）---")
    ppt_files = [f for f in os.listdir(PPT_DIR) if f.endswith(('.pptx', '.ppt', '.pdf'))]
    if ppt_files:
        print("发现以下PPT文件：")
        for i, f in enumerate(ppt_files):
            print(f"{i+1}. {f}")
        ppt_choice = input("请输入要使用的PPT编号（直接回车跳过）: ").strip()
        if ppt_choice:
            ppt_path = os.path.join(PPT_DIR, ppt_files[int(ppt_choice) - 1])
    else:
        print(f"请将PPT文件放入 '{PPT_DIR}' 文件夹后重试")
    
    if not os.path.exists(AUDIO_PATH):
        print(f"--- 提取音频中: {VIDEO_PATH} ---")
        extract_audio(VIDEO_PATH, AUDIO_PATH)
    else:
        print(f"--- 发现已有音频 {AUDIO_PATH}，跳过提取 ---")
    
    print(f"--- 正在加载模型识别视频: {target_video_name} ---")
    
    # 清理缓存并加载模型
    print("--- 正在清理缓存（解决可能的兼容性问题）---")
    import shutil
    try:
        # 清理可能损坏的缓存
        cache_dir = os.path.join(os.getcwd(), "cache")
        if os.path.exists(cache_dir):
            shutil.rmtree(cache_dir)
        
        # 清理 huggingface 缓存
        hf_cache = os.path.expanduser("~/.cache/huggingface")
        if os.path.exists(hf_cache):
            # 只清理 faster-whisper 相关的缓存
            hf_models = os.path.join(hf_cache, "models")
            if os.path.exists(hf_models):
                for item in os.listdir(hf_models):
                    if "Systran" in item or "whisper" in item.lower():
                        try:
                            shutil.rmtree(os.path.join(hf_models, item))
                            print(f"✅ 已清理: {item}")
                        except:
                            pass
    except Exception as e:
        print(f"⚠️ 清理缓存时出错: {e}")
    
    print("--- 正在加载 Whisper 模型 ---")
    # 根据配置选择使用 CPU 或 GPU
    device = "cpu" if USE_CPU else "cuda"
    compute_type = "int8" if USE_CPU else "float16"
    print(f"使用设备: {device}, 计算类型: {compute_type}")
    model = WhisperModel(MODEL_SIZE, device=device, compute_type=compute_type)
    segments = transcribe_audio(model, AUDIO_PATH)
    
    # 自动生成关键词/更新stop_words选项
    print("\n--- 智能优化选项（可选）---")
    print("1. 分析课程，生成专属关键词")
    print("2. 分析课程，更新stop_words（过滤无意义词汇）")
    print("3. 两者都做")
    print("0. 跳过（使用默认配置）")
    optimize_choice = input("请输入选项编号（默认0）: ").strip() or "0"
    
    final_keywords = DEFAULT_KEYWORDS
    
    if optimize_choice in ["1", "3"]:
        course_keywords = generate_course_keywords(segments, base_name)
        if course_keywords:
            final_keywords = list(set(DEFAULT_KEYWORDS + course_keywords))
            print(f"✅ 已添加课程专属关键词，当前关键词数: {len(final_keywords)}")
    
    if optimize_choice in ["2", "3"]:
        analyze_and_update_stop_words(segments)
    
    hits = find_keyword_hits(segments, final_keywords)
    
    if hits:
        # 先提取PPT文字（如果有）
        slides_text = None
        if ppt_path and os.path.exists(ppt_path) and ppt_path.endswith(('.pdf', '.pptx', '.ppt')):
            print("\n--- 正在提取PPT文字 ---")
            slides_text = extract_ppt_text(ppt_path)
            print(f"共提取 {len(slides_text)} 页PPT文字")
        
        # 分析每个命中点，提取相关内容并关联PPT
        raw_results = analyze_hits_and_extract_ppt(
            hits, segments, current_screenshot_dir, 
            ppt_path=ppt_path, slides_text=slides_text
        )
        
        # 过滤掉不相关的内容
        final_data = [r for r in raw_results if r.get('is_relevant', True) is not False]
        
        print(f"\n📊 原始命中: {len(raw_results)} 个 → 过滤后: {len(final_data)} 个相关内容")
        
        if final_data:
            # 收集所有相关的PPT页面
            all_ppt_pages = set()
            for item in final_data:
                if item.get('related_ppt_pages'):
                    all_ppt_pages.update(item['related_ppt_pages'])
            
            # 生成PDF
            if all_ppt_pages and ppt_path and ppt_path.endswith('.pdf'):
                # 提取相关PPT页面
                pdf_pages = [p - 1 for p in all_ppt_pages]
                temp_pdf = os.path.join(current_screenshot_dir, "selected_ppt.pdf")
                select_pdf_pages(ppt_path, pdf_pages, temp_pdf)
                ppt_selected_pdf = temp_pdf
                
                # 合并PPT和课堂笔记
                from pypdf import PdfReader, PdfWriter
                writer = PdfWriter()
                
                # 先添加PPT页面
                ppt_reader = PdfReader(ppt_selected_pdf)
                for page in ppt_reader.pages:
                    writer.add_page(page)
                
                # 添加课堂笔记
                temp_notes = os.path.join(RESULT_DIR, f"temp_notes_{os.getpid()}.pdf")
                create_final_pdf(final_data, f"{base_name} 课堂笔记", temp_notes)
                
                notes_reader = PdfReader(temp_notes)
                for page in notes_reader.pages:
                    writer.add_page(page)
            
                # 保存最终PDF
                with open(PDF_PATH, 'wb') as f:
                    writer.write(f)
                
                os.remove(temp_notes)
                print(f"✨ 已合并PPT和课堂笔记: {PDF_PATH}")
            elif all_ppt_pages and ppt_path and ppt_path.endswith(('.pptx', '.ppt')):
                # 将PPT相关页面转换为图片并生成PDF
                ppt_images = convert_ppt_to_images(ppt_path, list(all_ppt_pages), current_screenshot_dir)
                create_pdf_with_ppt(final_data, ppt_images, f"{base_name} 课堂笔记", PDF_PATH)
                print(f"✨ 处理完成！PDF 已生成到: {PDF_PATH}")
            else:
                create_final_pdf(final_data, f"{base_name} 课堂笔记", PDF_PATH)
                print(f"✨ 处理完成！PDF 已生成到: {PDF_PATH}")
        else:
            print("💡 没发现重点。")